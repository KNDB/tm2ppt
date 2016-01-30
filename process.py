import csv
import codecs
import locale
from collections import defaultdict
from pptx import Presentation, util

enc = locale.getpreferredencoding()
with codecs.open('ronde1.CSV', encoding=enc) as f:
	rankings = defaultdict(list)
	rankings_start = False
	for l in f.readlines(): 
		result = l.strip()
		if l.startswith('Ranglijst'):
			rankings_start = True
			rankings_nr = 1
			category = 'pup' if not rankings['pup'] else 'welp'
		elif rankings_start and result: 
			if result.startswith('AW'):
				continue
			if not result[0].isdigit(): 
				result = str(rankings_nr) + ';' + result
			rankings[category].append(result)
			rankings_nr += 1 
		
# print rankings
#	with codecs.open('process.csv', 'wb', 'utf-8') as p:
#		p.write('\r\n'.join(rankings))

prs = Presentation('Standen.pptx')

def fill_rows(table, category):
	for i, row in enumerate(table.rows):
		if i == 0: 
			continue
		ranking_cell = rankings[category][i-1].split(';')
		for j, cell in enumerate(row.cells): 
			if j == 0: 
				continue
			cell.text = ranking_cell[j]
			paragraph = cell.text_frame.paragraphs[0]
			paragraph.font.size = util.Pt(20)
			paragraph.font.name = 'Verdana'

fill_rows(prs.slides[2].shapes[2].table, 'pup')
fill_rows(prs.slides[1].shapes[2].table, 'welp')

prs.save('Standen-new.pptx')
